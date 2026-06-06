from pathlib import Path
import httpx
import chromadb
from chromadb.config import Settings as ChromaSettings
from typing import List, Tuple
import pypdf
import docx
from config import get_settings
from ollama_client import ollama_headers

settings = get_settings()

def get_chroma_client():
    client = chromadb.HttpClient(
        host=settings.chroma_host,
        port=settings.chroma_port,
        settings=ChromaSettings(anonymized_telemetry=False),
    )
    try: client.create_tenant(chromadb.DEFAULT_TENANT)
    except: pass
    try: client.create_database(chromadb.DEFAULT_DATABASE)
    except: pass
    return client

async def get_embeddings(texts: List[str]) -> List[List[float]]:
    async with httpx.AsyncClient(
        timeout=httpx.Timeout(connect=30, read=300, write=60, pool=30), verify=True
    ) as client:
        resp = await client.post(
            f"{settings.ollama_base_url}/api/embed",
            headers=ollama_headers(),
            json={"model": "nomic-embed-text", "input": texts},
        )
        resp.raise_for_status()
        data = resp.json()
        return data.get("embeddings") or [data.get("embedding")]

def extract_text_from_pdf(path: str) -> str:
    reader = pypdf.PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def extract_text_from_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_text_from_xlsx(path: str) -> str:
    ext = Path(path).suffix.lower()
    lines = []
    if ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(path)
        for sheet in wb.sheets():
            lines.append(f"=== Feuille : {sheet.name} ===")
            for row_idx in range(sheet.nrows):
                cells = [str(sheet.cell_value(row_idx, c))
                         for c in range(sheet.ncols)
                         if sheet.cell_value(row_idx, c) != ""]
                if cells:
                    lines.append(" | ".join(cells))
    else:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            lines.append(f"=== Feuille : {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(" | ".join(cells))
    return "\n".join(lines)

def extract_text_from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"=== Diapositive {i+1} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)

def extract_text_from_txt(path: str) -> str:
    import chardet
    with open(path, "rb") as f:
        raw = f.read()
    enc = chardet.detect(raw)["encoding"] or "utf-8"
    return raw.decode(enc, errors="replace")

def extract_text(path: str, file_type: str) -> str:
    ext = file_type.lower()
    if ext == "pdf":            return extract_text_from_pdf(path)
    elif ext in ("docx","doc"): return extract_text_from_docx(path)
    elif ext in ("xlsx","xls"): return extract_text_from_xlsx(path)
    elif ext in ("pptx","ppt"): return extract_text_from_pptx(path)
    else:                       return extract_text_from_txt(path)

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i: i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        i += chunk_size - overlap
    return chunks

async def index_document(file_path, file_type, user_id, doc_id) -> Tuple[str, int]:
    text = extract_text(file_path, file_type)
    if not text.strip(): return "", 0
    chunks = chunk_text(text)
    if not chunks: return "", 0
    all_embeddings = []
    for i in range(0, len(chunks), 20):
        embs = await get_embeddings(chunks[i:i+20])
        all_embeddings.extend(embs)
    collection_name = f"user_{user_id}_doc_{doc_id}"
    client = get_chroma_client()
    col = client.get_or_create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})
    col.add(documents=chunks, embeddings=all_embeddings, ids=[f"chunk_{i}" for i in range(len(chunks))])
    return collection_name, len(chunks)

async def query_document(collection_name, question, n_results=4) -> List[str]:
    q_embed = await get_embeddings([question])
    try:
        col = get_chroma_client().get_collection(collection_name)
        results = col.query(query_embeddings=q_embed, n_results=n_results)
        return results["documents"][0] if results["documents"] else []
    except: return []

async def query_user_documents(user_id, collection_names, question, n_results=3) -> str:
    if not collection_names: return ""
    all_chunks = []
    for col in collection_names:
        all_chunks.extend(await query_document(col, question, n_results))
    if not all_chunks: return ""
    return "Contexte extrait des documents :\n\n" + "\n\n---\n\n".join(all_chunks[:6]) + "\n\n"

async def delete_document_collection(collection_name):
    try: get_chroma_client().delete_collection(collection_name)
    except: pass
